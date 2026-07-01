# -*- coding: utf-8 -*-
"""
요양(병)원 시설 유형별 대피체계 표준모델 개발 - 요약제안서 PPT 생성
디자인: 심플하고 깔끔한 스타일
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────────────────────────
# 디자인 상수
# ──────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2E, 0x86, 0xAB)
ACCENT    = RGBColor(0x3D, 0x5A, 0x80)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x33, 0x33, 0x33)
GRAY      = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
TBL_HEADER = RGBColor(0x1B, 0x2A, 0x4A)
TBL_ROW1   = RGBColor(0xF7, 0xF9, 0xFC)
TBL_ROW2   = RGBColor(0xE8, 0xEE, 0xF4)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_TITLE = "맑은 고딕"
FONT_BODY  = "맑은 고딕"


# ──────────────────────────────────────────────
# 유틸리티 함수
# ──────────────────────────────────────────────
def add_bottom_bar(slide):
    """슬라이드 하단에 네이비 바 추가"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.1),
        SLIDE_WIDTH, Inches(0.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()


def add_top_accent_line(slide):
    """슬라이드 상단에 얇은 악센트 라인"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = DARK_BLUE
    line.line.fill.background()


def add_page_number(slide, num, total):
    """페이지 번호 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(12.2), Inches(7.15), Inches(1), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.RIGHT


def add_section_header(slide, section_text):
    """섹션 배지 (좌상단)"""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.25),
        Inches(3.5), Inches(0.45)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = NAVY
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.name = FONT_BODY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)


def set_slide_bg(slide, color=WHITE):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, left, top, width, height, font_size=28, color=NAVY, bold=True, alignment=PP_ALIGN.LEFT):
    """제목 텍스트 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_TITLE
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, font_size=13, color=BLACK, bold=False, line_spacing=1.5):
    """본문 텍스트 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_BODY
    p.font.bold = bold
    p.space_after = Pt(4)
    p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_bullet_list(slide, items, left, top, width, height, font_size=12, color=BLACK, bullet_color=DARK_BLUE):
    """글머리 기호 리스트 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 불릿 문자 + 텍스트
        run_bullet = p.add_run()
        run_bullet.text = "●  "
        run_bullet.font.size = Pt(7)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = FONT_BODY
        
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = FONT_BODY
        
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * 1.4)
    return tf


def add_sub_bullet_list(slide, items, left, top, width, height, font_size=11, color=GRAY):
    """하위 글머리 리스트"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        run_bullet = p.add_run()
        run_bullet.text = "    ─  "
        run_bullet.font.size = Pt(8)
        run_bullet.font.color.rgb = LIGHT_GRAY
        run_bullet.font.name = FONT_BODY
        
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = FONT_BODY
        
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * 1.3)
    return tf


def add_table(slide, rows, cols, data, left, top, width, height, col_widths=None):
    """테이블 추가 (헤더행 네이비 배경)"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 열 너비 설정
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ""
            
            # 셀 텍스트 설정
            p = cell.text_frame.paragraphs[0]
            p.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 셀 마진
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            
            if r == 0:  # 헤더
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_HEADER
            else:
                p.font.size = Pt(10)
                p.font.color.rgb = BLACK
                p.font.bold = False
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_ROW1 if r % 2 == 1 else TBL_ROW2
    
    return table


def add_highlight_box(slide, text, left, top, width, height, bg_color=LIGHT_BG, text_color=NAVY, font_size=13, bold=True):
    """강조 박스 추가"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = DARK_BLUE
    box.line.width = Pt(1)
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = FONT_BODY
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return tf


def add_divider_line(slide, left, top, width):
    """구분선 추가"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()


# ──────────────────────────────────────────────
# PPT 생성 시작
# ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

BLANK_LAYOUT = prs.slide_layouts[6]  # blank layout
TOTAL_SLIDES = 25


# ============================================================
# 슬라이드 1: 표지
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, NAVY)

# 상단 악센트 라인
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = DARK_BLUE
line.line.fill.background()

# 중앙 타이틀
add_title_text(slide, "요양(병)원 시설 유형별\n대피체계 표준모델 개발",
               Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.5),
               font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 구분선
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.03))
divider.fill.solid()
divider.fill.fore_color.rgb = DARK_BLUE
divider.line.fill.background()

# 부제목
add_title_text(slide, "요  약  제  안  서",
               Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.7),
               font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8), bold=False, alignment=PP_ALIGN.CENTER)

# 날짜
add_title_text(slide, "2026년",
               Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
               font_size=16, color=RGBColor(0x88, 0xAA, 0xCC), bold=False, alignment=PP_ALIGN.CENTER)


# ============================================================
# 슬라이드 2: 목차
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 2, TOTAL_SLIDES)

add_title_text(slide, "목  차", Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
               font_size=32, color=NAVY)

add_divider_line(slide, Inches(0.8), Inches(1.2), Inches(11.7))

# 목차 항목들
toc_sections = [
    ("Ⅰ", "제안 개요 및 추진 전략", [
        "추진 배경 및 문제 제기",
        "과업 목표 및 수행 범위",
        "핵심 역량 및 제안의 차별화",
        "제안기관의 핵심 역량"
    ]),
    ("Ⅱ", "세부 연구 수행 계획", [
        "과업1: 기초현황조사 (시설 분류·사례 분석·심층 인터뷰)",
        "과업2: 국내외 대피체계 비교 분석",
        "과업3: 대피체계 표준모델 (정량적 근거·표준안·행동 매뉴얼)",
        "과업4: 합동 훈련 시나리오 및 정책 과제",
        "중장기 확산 로드맵"
    ]),
    ("Ⅲ", "사업 관리 및 기관 역량", [
        "연구 추진 일정 및 마일스톤",
        "보고 체계 및 품질관리",
        "제안기관 핵심 역량"
    ])
]

y_pos = 1.5
for sec_num, sec_title, sub_items in toc_sections:
    # 섹션 번호 + 제목
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y_pos), Inches(0.5), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = NAVY
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = sec_num
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = FONT_BODY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(2)
    
    add_title_text(slide, sec_title, Inches(1.7), Inches(y_pos), Inches(5), Inches(0.4),
                   font_size=16, color=NAVY, bold=True)
    
    y_pos += 0.5
    for item in sub_items:
        add_body_text(slide, f"    {item}", Inches(2.0), Inches(y_pos), Inches(9), Inches(0.3),
                      font_size=12, color=GRAY)
        y_pos += 0.32
    y_pos += 0.2


# ============================================================
# 슬라이드 3: 추진 배경 및 문제 제기
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 3, TOTAL_SLIDES)
add_section_header(slide, "Ⅰ. 제안 개요 및 추진 전략")

add_title_text(slide, "1. 추진 배경 및 문제 제기",
               Inches(0.8), Inches(0.85), Inches(8), Inches(0.6), font_size=24, color=NAVY)

add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 좌측: 초고령사회
add_title_text(slide, "초고령사회 진입과 요양병원 급증",
               Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=15, color=ACCENT, bold=True)

add_bullet_list(slide, [
    "2004년 92개 → 2021년 1,466개로 요양병원 대폭 증가",
    "입소 환자 약 90%가 자력대피 불가능한 '와상환자'",
    "야간 시간대 의료인 1명당 과도한 환자수 부담",
], Inches(0.8), Inches(2.15), Inches(5.5), Inches(2.0), font_size=12)

# 좌측: 구조적 한계
add_title_text(slide, "기존 대피체계의 구조적 한계",
               Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.4), font_size=15, color=ACCENT, bold=True)

add_bullet_list(slide, [
    "일률적 대피체계 → 시설 유형별 현장 적용성 저하",
    "밀양 세종병원(2018): 사망 47명, 방화구획 미비",
    "장성 효사랑(2014): 사망 21명, 조력인력 3명 불과",
], Inches(0.8), Inches(4.25), Inches(5.5), Inches(1.8), font_size=12)

# 우측: RSET > ASET 핵심 문제 강조 박스
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xF0)
box.line.color.rgb = RED_ACCENT
box.line.width = Pt(1.5)

add_title_text(slide, "⚠  핵심 문제: RSET > ASET",
               Inches(7.3), Inches(1.9), Inches(5), Inches(0.5),
               font_size=18, color=RED_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_body_text(slide, "\"현 상태로는 피난 자체가 불가능\"",
              Inches(7.3), Inches(2.5), Inches(5), Inches(0.4),
              font_size=14, color=RED_ACCENT, bold=True)

# ASET/RSET 테이블
tbl_data = [
    ["구분", "수치"],
    ["피난허용시간(ASET)", "53.22초"],
    ["피난소요시간(RSET)", "76.80초"],
    ["판정", "피난안전성 미확보"],
]
add_table(slide, 4, 2, tbl_data,
          Inches(7.5), Inches(3.2), Inches(4.8), Inches(2.2),
          col_widths=[2.4, 2.4])

add_body_text(slide, "출처: 박국희(2023) 피난 시뮬레이션 연구",
              Inches(7.5), Inches(5.5), Inches(4.8), Inches(0.3),
              font_size=9, color=GRAY)


# ============================================================
# 슬라이드 4: 과업 목표 및 수행 범위
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 4, TOTAL_SLIDES)
add_section_header(slide, "Ⅰ. 제안 개요 및 추진 전략")

add_title_text(slide, "2. 과업 목표 및 수행 범위",
               Inches(0.8), Inches(0.85), Inches(8), Inches(0.6), font_size=24, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 3대 핵심 목표
add_title_text(slide, "3대 핵심 목표",
               Inches(0.8), Inches(1.65), Inches(5), Inches(0.4), font_size=16, color=ACCENT, bold=True)

tbl_data = [
    ["구분", "목표", "핵심 전략"],
    ["목표 1", "시설 유형별 위험요인 정량화\n및 맞춤형 대피체계 구축", "입지·건물·설비·환자 중증도별\n세분화된 위험 프로파일 도출"],
    ["목표 2", "와상환자 중심의\n'수평피난(Defend-in-Place)' 반영", "무리한 수직 이동 대신,\n층내 방화구획 안전구역으로 수평이동"],
    ["목표 3", "행안부 기존 가이드라인의\n현장 작동형 매뉴얼 고도화", "법정 소방계획서 양식 연계,\n종사자 행정 부담 최소화"],
]
add_table(slide, 4, 3, tbl_data,
          Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.8),
          col_widths=[1.5, 5.0, 5.2])

# 4대 주요 산출물
add_title_text(slide, "4대 주요 산출물",
               Inches(0.8), Inches(5.2), Inches(5), Inches(0.4), font_size=16, color=ACCENT, bold=True)

tbl_data2 = [
    ["산출물", "내용"],
    ["유형별 대피계획 표준안", "건축·설비적 대안 포함, 소방계획서 연계형 표준 서식"],
    ["행동 매뉴얼", "환자 중증도별(A/B/C유형) 맞춤 대피동선 + 종사자 4단계 임무 규격화"],
    ["합동 훈련 시나리오", "도상/현장모의/전문가참여 3종, 주간·야간 분리, 타임라인 기반"],
    ["정책 과제", "법령 개정안, A~D등급 자가진단 도구, 우수시설 인센티브 제도"],
]
add_table(slide, 5, 2, tbl_data2,
          Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.3),
          col_widths=[3.0, 8.7])


# ============================================================
# 슬라이드 5: 차별화 포인트
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 5, TOTAL_SLIDES)
add_section_header(slide, "Ⅰ. 제안 개요 및 추진 전략")

add_title_text(slide, "3. 핵심 역량 및 제안의 차별화",
               Inches(0.8), Inches(0.85), Inches(8), Inches(0.6), font_size=24, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 3개 차별화 포인트를 박스로
box_data = [
    ("시뮬레이션 비용·시간 효율화", [
        "신규 시뮬레이션 용역 발주 없이 학계 교차 검증된 정량적 데이터 활용",
        "제안사 보유 FDS 자동화 특허(국립재난안전연구원 공동발명) 적용",
        "분석 정밀도 및 속도 극대화"
    ]),
    ("현장 중심 실효성 확보", [
        "독립 매뉴얼 추가 배포 방식 지양 → 행정 부담 최소화",
        "기존 법정 '소방계획서' 양식과 구조적 연계",
        "현장 수용성과 실무 적용성 극대화"
    ]),
    ("환자 3유형 맞춤 동선 설계", [
        "A유형(보행 가능): 핸드레일·바닥 발광 유도선 활용 자력 대피",
        "B유형(휠체어 부분도움): 종사자 1인 전담 보조, 경사로 활용",
        "C유형(거동불가 와상): 수직 피난 금지, 2인 이상 팀 편성 수평 대피"
    ]),
]

x_positions = [0.8, 4.8, 8.8]
for i, (title, items) in enumerate(box_data):
    x = x_positions[i]
    # 박스 배경
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.7), Inches(3.7), Inches(5.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = RGBColor(0xD0, 0xDA, 0xE8)
    box.line.width = Pt(0.75)
    
    # 번호 배지
    num_badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x + 1.45), Inches(1.9), Inches(0.45), Inches(0.45))
    num_badge.fill.solid()
    num_badge.fill.fore_color.rgb = DARK_BLUE
    num_badge.line.fill.background()
    tf = num_badge.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(2)
    
    # 제목
    add_title_text(slide, title, Inches(x + 0.2), Inches(2.5), Inches(3.3), Inches(0.5),
                   font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 구분선
    add_divider_line(slide, Inches(x + 0.3), Inches(3.05), Inches(3.1))
    
    # 내용
    add_bullet_list(slide, items, Inches(x + 0.2), Inches(3.2), Inches(3.3), Inches(3.3),
                    font_size=10, bullet_color=DARK_BLUE)


# ============================================================
# 슬라이드 6: 제안기관 핵심 역량
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 6, TOTAL_SLIDES)
add_section_header(slide, "Ⅰ. 제안 개요 및 추진 전략")

add_title_text(slide, "4. 제안기관의 핵심 역량",
               Inches(0.8), Inches(0.85), Inches(8), Inches(0.6), font_size=24, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 특허 정보
add_title_text(slide, "행안부 국립재난안전연구원 공동 특허 발명",
               Inches(0.8), Inches(1.7), Inches(8), Inches(0.4), font_size=16, color=ACCENT, bold=True)

tbl_data = [
    ["항목", "내용"],
    ["발명 명칭", "GUI 기반의 화재 시뮬레이션 입력 데이터 자동 생성 방법 및 장치"],
    ["출원번호", "10-2026-0080248 (2026.05.04)"],
    ["출원인", "대한민국(행정안전부 국립재난안전연구원)"],
    ["공동 발명자", "국립재난안전연구원 + 제안사"],
]
add_table(slide, 5, 2, tbl_data,
          Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.0),
          col_widths=[2.5, 9.2])

# 국책 과제
add_title_text(slide, "대형 국책 과제 수행 실적",
               Inches(0.8), Inches(4.5), Inches(8), Inches(0.4), font_size=16, color=ACCENT, bold=True)

tbl_data2 = [
    ["사업명", "발주처", "기간"],
    ["기후변화 대응 재해예방사업 표준화 기술개발", "행안부", "2021.04~2025.12"],
    ["도시물순환 연계 홍수 취약성 진단 및 효과 분석 연구", "한국수자원공사", "2023.07~2024.12"],
    ["댐운영 의사결정 고도화 기본구상 수립 용역", "한국수자원공사", "2025.12~현재"],
    ["공주시 통합물관리 기본계획 수립 용역", "충남 공주시", "2021.08~2025.07"],
]
add_table(slide, 5, 3, tbl_data2,
          Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8),
          col_widths=[7.0, 2.5, 2.2])


# ============================================================
# 슬라이드 7: 과업1 - 시설 유형 분류 매트릭스
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 7, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업1. 기초현황조사 ─ 시설 유형 분류 매트릭스",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

tbl_data = [
    ["분류", "세부 유형", "위험요인 분석 방향"],
    ["입지", "도심형, 교외형,\n산림인접형, 농어촌형", "도심형: 외부 대피공간 확보 곤란\n산림인접형: 외부 산불 연속 확대"],
    ["건물 유형", "단일용도, 복합용도,\n고층형, 저층분산형", "복합용도: 저층부 타 시설 화재 전이\n고층형: 수직피난 다중밀집"],
    ["준공 시기", "스프링클러 의무화\n전(과도기)/후(완비)", "미설치: 플래쉬오버 단시간 발생,\nASET 급격 단축"],
    ["환자 중증도", "A유형(보행)/B유형(휠체어)\n/C유형(거동불가)", "C유형 80% 이상 시\n'최고 위험군' 분류"],
]
add_table(slide, 5, 3, tbl_data,
          Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.2),
          col_widths=[2.0, 4.0, 5.7])

# 위험 프로파일 설명
add_title_text(slide, "유형 조합별 위험 프로파일(Risk Profile) 도출",
               Inches(0.8), Inches(5.2), Inches(8), Inches(0.4), font_size=15, color=ACCENT, bold=True)

add_bullet_list(slide, [
    "4가지 분류축의 세부 유형을 교차 조합하여 시설 고유의 위험 프로파일 산출",
    "예시: [도심형 + 복합용도 + 고층형 + 미설치 + C유형 80% 이상] → 최고 위험군 지정",
    "작성된 위험 프로파일은 '유형별 맞춤형 대피계획 및 행동 매뉴얼' 설계의 직접적 기초 데이터로 활용",
], Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.3), font_size=11)


# ============================================================
# 슬라이드 8: 과업1 - 화재 대피 사례 분석
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 8, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업1. 기초현황조사 ─ 화재 대피 사례 심층 분석",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

add_title_text(slide, "실패 요인 도출",
               Inches(0.8), Inches(1.65), Inches(5), Inches(0.4), font_size=16, color=ACCENT, bold=True)

tbl_data = [
    ["핵심 분석 항목", "도출된 문제점"],
    ["대피 인원 분류 부재", "환자 신체기능별 피난계획 없이 일률적 동선 적용 → 피난구 밀집·혼란"],
    ["조력 인력의 물리적 한계", "C유형 1명 대피에 최소 2인 필요, 심야 참사 시 인력 절대 부족"],
    ["ASET 초과 현상", "방화구획+가연물 연소 → ASET 단축, RSET 증가 → 골든타임 실패"],
    ["수평 피난 수단 부재", "수직 대피 의존, 대피공간·경사로·연결복도 등 수평 대피 인프라 미비"],
]
add_table(slide, 5, 2, tbl_data,
          Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.5),
          col_widths=[3.5, 8.2])

# 하단 심층 인터뷰 요약
add_title_text(slide, "현장 중심 심층 인터뷰 계획",
               Inches(0.8), Inches(5.0), Inches(5), Inches(0.4), font_size=16, color=ACCENT, bold=True)

add_bullet_list(slide, [
    "조사 대상: 시설 운영자·소방안전관리자, 요양보호사·의료진, 119구급대원",
    "표본: 시설 유형 분류 매트릭스 기준 최소 10개소 이상 현장",
    "조사 방법: 핵심 질문 기반 반구조화 면접 (자유 의견 개진 유도)",
], Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.2), font_size=11)


# ============================================================
# 슬라이드 9: 과업1 - 3대 중점 조사 항목
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 9, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업1. 현장 심층 인터뷰 ─ 3대 중점 조사 항목",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

items_data = [
    ("① 야간 시간대 인력 현황 및 조력 피난 한계", [
        "대부분의 대형 인명피해는 최소 인력 근무 심야 시간대에 발생",
        "규정상 야간 인력과 C유형 환자 대피 최소 필요 인력(2인 이상) 간의 물리적 격차 조사",
        "수직 피난 불가 원인 규명 → 수평 피난 및 Defend-in-Place 전략 도입 근거 확보",
    ]),
    ("② 대피 훈련 실태 및 매뉴얼 작동 여부", [
        "현행 소방 훈련이 서류상 요건 충족에 그치는지 실태 점검",
        "노인 대피 행동 특성(추종본능, 지광본능, 익숙한 공간 집착) 인지 여부 파악",
        "피난 기구(휠체어, 수평/수직 대피 도움장치 등) 실제 사용 경험 수렴",
    ]),
    ("③ 현장 대응 애로사항 및 유관기관 협조 체계", [
        "수직 피난 시 계단 폭 협소로 인한 다중밀집(병목) 체감도",
        "119구급대 도착 시 환자 인계 과정 지연 요인 파악",
        "치매환자 돌발 행동 통제 어려움 → 합동 훈련 시나리오 중점 개선 과제로 연계",
    ]),
]

y = 1.7
for title, bullets in items_data:
    add_title_text(slide, title, Inches(0.8), Inches(y), Inches(11), Inches(0.4),
                   font_size=14, color=NAVY, bold=True)
    y += 0.45
    add_bullet_list(slide, bullets, Inches(1.2), Inches(y), Inches(11), Inches(1.2),
                    font_size=11, bullet_color=DARK_BLUE)
    y += 1.35


# ============================================================
# 슬라이드 10: 과업2 - 국내외 대피체계 비교
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 10, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업2. 국내외 대피체계 비교 분석",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

add_title_text(slide, "국가별 피난 및 안전관리 비교 매트릭스",
               Inches(0.8), Inches(1.65), Inches(8), Inches(0.4), font_size=15, color=ACCENT, bold=True)

tbl_data = [
    ["비교 항목", "한국(현행)", "미국(NFPA 101)", "일본(건축/소방법)"],
    ["대피 전략", "전관·수직 피난 중심", "제자리 방어\n(Defend-in-Place) 원칙", "수평구획 및\n일시대기공간 중심"],
    ["시설 기준", "대피공간·경사로 등\n피난로 확보 의무화", "철저한 방화구획 및\n층별 안전구역 의무화", "면적 제한 없는\n소방설비 소급 적용"],
    ["인력 기준", "법정 최소 상주 인원\n(야간 취약성 잔존)", "훈련된 전담\n감시인·조력자 상시 배치", "인력 한계 극복 위한\n설비 자동화 연동"],
    ["협업 체계", "화재 발생 후\n소방·의료기관 개입", "지역사회 기반\n사전 대응 프로토콜", "소방기관과 자동 연동된\n즉각적 통보체계"],
]
add_table(slide, 5, 4, tbl_data,
          Inches(0.8), Inches(2.1), Inches(11.7), Inches(3.5),
          col_widths=[2.0, 3.2, 3.2, 3.3])

# 하단 시사점
add_highlight_box(slide,
    "시사점: '제자리 방어' 기반 → 수직 피난 중심에서 수평 피난 전략 패러다임 전환 필요",
    Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
    font_size=13, bold=True)


# ============================================================
# 슬라이드 11: 과업2 - 미국/일본 핵심 분석
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 11, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업2. 해외 제도 핵심 분석 및 한국형 적용 방안",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 미국 NFPA
add_title_text(slide, "미국 NFPA 101 핵심 분석",
               Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=15, color=ACCENT, bold=True)
add_bullet_list(slide, [
    "요양병원을 '지체부자유자 보호시설'로 분류",
    "수평피난: 1개 층을 2개 이상 방화구획으로 분할",
    "환자를 비화재존으로 침상 채 이동",
    "3층 초과 건물: 최소 2시간 내화성능 확보",
], Inches(0.8), Inches(2.15), Inches(5.5), Inches(2.0), font_size=11)

# 일본
add_title_text(slide, "일본 고령자 시설 화재 대응 제도",
               Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4), font_size=15, color=ACCENT, bold=True)
add_bullet_list(slide, [
    "스프링클러 전면 의무화 (면적 제한 철폐)",
    "거실·복도·계단실 철저한 방화구획",
    "외기 개방 발코니 등 일시대기공간 활용",
    "소규모 시설 포함 예외 없는 적용",
], Inches(7.0), Inches(2.15), Inches(5.5), Inches(2.0), font_size=11)

# 한국형 적용
add_title_text(slide, "한국형 적용 방안 ─ 소방법 소급 적용",
               Inches(0.8), Inches(4.3), Inches(10), Inches(0.4), font_size=15, color=ACCENT, bold=True)

tbl_data = [
    ["설비 항목", "적용 기준"],
    ["스프링클러", "바닥면적 합계 600㎡ 이상은 모든 층에 설치"],
    ["간이스프링클러", "면적 관계없이 전체 노인복지시설 설치 의무화"],
    ["자동화재탐지설비", "면적 관계없이 전체 노인복지시설 설치 의무화"],
    ["자동화재속보설비", "면적 관계없이 전체 노인복지시설 설치 의무화"],
    ["직통계단", "3층 이상 노유자시설, 바닥면적 200㎡ 이상 시 2개소 이상"],
    ["배연설비", "6층 이상 또는 노인요양시설은 층수 관계없이 설치"],
]
add_table(slide, 7, 2, tbl_data,
          Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0),
          col_widths=[3.0, 8.7])


# ============================================================
# 슬라이드 12: 과업3 - 이중 접근 전략
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 12, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업3. 대피체계 표준모델 ─ 기술적·정량적 근거",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 핵심 공식 강조
add_highlight_box(slide,
    "피난안전성 확보 조건:   ASET(피난허용시간)  >  RSET(피난소요시간)",
    Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.55),
    bg_color=RGBColor(0xE8, 0xF4, 0xFD), text_color=NAVY, font_size=16, bold=True)

# 이중 접근 전략
add_title_text(slide, "이중 접근(Two-Track) 전략",
               Inches(0.8), Inches(2.55), Inches(8), Inches(0.4), font_size=16, color=ACCENT, bold=True)

# RSET 단축 박스
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(3.0), Inches(5.6), Inches(1.2))
box1.fill.solid()
box1.fill.fore_color.rgb = LIGHT_BG
box1.line.color.rgb = DARK_BLUE
box1.line.width = Pt(1)

add_title_text(slide, "RSET 단축 (피난소요시간 줄이기)",
               Inches(1.0), Inches(3.1), Inches(5.2), Inches(0.35),
               font_size=13, color=DARK_BLUE, bold=True)
add_body_text(slide, "수평대피공간 확보, 복합 피난수단 동시 운용,\n종사자 임무 명확화",
              Inches(1.0), Inches(3.5), Inches(5.2), Inches(0.6),
              font_size=11, color=BLACK)

# ASET 연장 박스
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.9), Inches(3.0), Inches(5.6), Inches(1.2))
box2.fill.solid()
box2.fill.fore_color.rgb = LIGHT_BG
box2.line.color.rgb = DARK_BLUE
box2.line.width = Pt(1)

add_title_text(slide, "ASET 연장 (피난허용시간 늘리기)",
               Inches(7.1), Inches(3.1), Inches(5.2), Inches(0.35),
               font_size=13, color=DARK_BLUE, bold=True)
add_body_text(slide, "연기감지기 교체, 배연창 연동 설치,\n내부 마감재 난연·불연화",
              Inches(7.1), Inches(3.5), Inches(5.2), Inches(0.6),
              font_size=11, color=BLACK)

# 4대 핵심 문제 테이블
add_title_text(slide, "4대 핵심 문제 → 기술적 대안 → 정량적 기대효과",
               Inches(0.8), Inches(4.5), Inches(10), Inches(0.4), font_size=15, color=ACCENT, bold=True)

tbl_data = [
    ["문제점", "기술적 대안", "정량적 기대효과", "근거"],
    ["ASET < RSET\n(피난불가)", "수평 대피공간 확보", "RSET 33.6% 단축\n(651.0초→432.0초)", "박국희\n(2023)"],
    ["열감지기\n감지 지연", "연기감지기 교체", "초기 화재 감지시간\n170초 단축", "배태훈\n(2018)"],
    ["연기 확산에\n따른 고립", "배연창 설치\n(감지기 연동)", "ASET 88.33초 증가", "배태훈\n(2018)"],
    ["단일 피난로\n병목 현상", "복합 피난수단\n동시 확보", "RSET 최대 46.37% 단축\n(805.5초→432.0초)", "박국희\n(2023)"],
]
add_table(slide, 5, 4, tbl_data,
          Inches(0.8), Inches(4.95), Inches(11.7), Inches(2.0),
          col_widths=[2.2, 2.8, 3.7, 3.0])


# ============================================================
# 슬라이드 13: 시나리오별 RSET 비교
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 13, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업3. 시나리오별 RSET 비교 분석",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

add_body_text(slide, "※ 경사로, 승강기, 대피공간을 모두 확보해야만 물리적 피난 한계를 극복할 수 있음 (박국희, 2023)",
              Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
              font_size=12, color=GRAY, bold=False)

tbl_data = [
    ["시나리오", "피난수단", "RSET", "시나리오1 대비 단축"],
    ["1", "경사로만", "805.5초", "기준 (0%)"],
    ["2", "승강기만", "686.8초", "14.74% ↓"],
    ["3", "경사로 + 승강기", "651.0초", "19.18% ↓"],
    ["4", "승강기 + 대피공간", "561.2초", "30.33% ↓"],
    ["5", "경사로 + 대피공간", "540.8초", "32.86% ↓"],
    ["6", "경사로 + 승강기 + 대피공간", "432.0초", "46.37% ↓"],
]
add_table(slide, 7, 4, tbl_data,
          Inches(0.8), Inches(2.1), Inches(11.7), Inches(3.2),
          col_widths=[1.5, 4.0, 2.5, 3.7])

# 핵심 시사점
add_highlight_box(slide,
    "핵심: 시나리오 6 (복합 피난수단 모두 확보) → RSET 46.37% 단축, 805.5초 → 432.0초",
    Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.55),
    bg_color=RGBColor(0xE8, 0xF4, 0xFD), text_color=NAVY, font_size=13, bold=True)


# ============================================================
# 슬라이드 14: 유형별 대피계획 표준안
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 14, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업3. 표준안 ─ 유형별 대피계획",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 대피 원칙 테이블
tbl_data = [
    ["항목", "내용"],
    ["대피 우선순위", "C유형(거동불가 와상환자) → B유형(휠체어) → A유형(보행가능)"],
    ["대피 원칙", "동일 층 인접 방화구획(비화재존)으로 '침상 채 수평 이동' 우선"],
    ["상황별\n대피 장소", "주간: 1차 수평대피 → 상황에 따라 수직 피난\n야간·심야: 수직 피난 배제, 층내 대피공간 → 소방대 구조 대기"],
    ["인력배치표", "교대 근무별(주간/야간/심야) 구역·임무 담당자 실명 기재 매트릭스"],
]
add_table(slide, 5, 2, tbl_data,
          Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.5),
          col_widths=[2.5, 9.2])

# 시설 유형별 대피계획서 작성례
add_title_text(slide, "시설 유형별 대피계획서 작성례 (최소 3개 유형)",
               Inches(0.8), Inches(4.5), Inches(8), Inches(0.4), font_size=15, color=ACCENT, bold=True)

tbl_data2 = [
    ["시설 유형", "대피 계획"],
    ["도심 고층형", "층별 2개 이상 방화구획 분할, 층내 수평 대피, 피난용 승강기 이용 통제"],
    ["교외 저층형/단독형", "경사로·넓은 출입구 활용, 휠체어·침대 환자 신속 옥외 대피 동선"],
    ["산림인접형", "반대편 옥외 대피소 지정, 구급차·시설차량 조기 외부 이송"],
]
add_table(slide, 4, 2, tbl_data2,
          Inches(0.8), Inches(4.95), Inches(11.7), Inches(1.8),
          col_widths=[2.5, 9.2])


# ============================================================
# 슬라이드 15: 종사자 4단계 행동 요령
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 15, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업3. 표준안 ─ 종사자 4단계 행동 요령",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

tbl_data = [
    ["단계", "역할", "주요 임무 행동"],
    ["1단계\n발견자", "최초 발견 직원", "\"불이야\" 육성 전파,\n소화기 초기 진압 시도 → 천장 번짐 시 즉시 대피 전환"],
    ["2단계\n전파자", "상황본부 직원", "비상벨 작동, 원내 화재 방송(Code Red),\n119 즉각 신고"],
    ["3단계\n대피지원자", "인력배치표 기준", "C → B → A유형 순서\n물리적 이송·유도 수행"],
    ["4단계\n확인자", "최종 점검 직원", "빈 방 출입문·방화문 확실 폐쇄"],
]
add_table(slide, 5, 3, tbl_data,
          Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.8),
          col_widths=[2.0, 3.0, 6.7])


# ============================================================
# 슬라이드 16: 야간 특화 시나리오
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 16, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업3. 야간(최소 인력) 시간대 특화 시나리오",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

scenarios = [
    ("1안", "단거리 수평 피난 모델 (구조 대기 집중)",
     "무리한 수직 피난을 배제하고, 와상환자를 동일 층 인접 방화구획으로 우선 이동",
     "방화문 완벽 차단 → 연기 유입 차단 → 소방대 도착 전 골든타임 확보"),
    ("2안", "중증도 기반 사전 재배치 모델 (위험 원천 제거)",
     "고위험 환자군을 1층 또는 피난층 등 대피가 용이한 구역에 선제적으로 배치",
     "심야 시간대 침대 이동 동선 자체를 원천적으로 최소화"),
    ("3안", "근무자 부담 분산 모델 (자력 피난 환경 조성)",
     "보행 가능 환자(A유형)가 시각적 유도선으로 스스로 옥외 대피할 수 있는 환경 구축",
     "한정된 야간 근무자가 C유형 구조에만 전념 가능"),
]

y = 1.7
for num, title, strategy, effect in scenarios:
    # 번호 배지
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(y), Inches(0.6), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(2)
    
    add_title_text(slide, title, Inches(1.6), Inches(y), Inches(10), Inches(0.4),
                   font_size=14, color=NAVY, bold=True)
    
    add_body_text(slide, f"전략: {strategy}",
                  Inches(1.6), Inches(y + 0.45), Inches(10.5), Inches(0.4),
                  font_size=11, color=BLACK)
    add_body_text(slide, f"기대효과: {effect}",
                  Inches(1.6), Inches(y + 0.85), Inches(10.5), Inches(0.4),
                  font_size=11, color=DARK_BLUE, bold=True)
    
    if num != "3안":
        add_divider_line(slide, Inches(0.8), Inches(y + 1.35), Inches(11.7))
    
    y += 1.55


# ============================================================
# 슬라이드 17: 과업4 - 합동 훈련 시나리오
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 17, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업4. 합동 훈련 시나리오",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 유관기관 협업 체계
add_title_text(slide, "유관기관 협업 체계",
               Inches(0.8), Inches(1.65), Inches(5), Inches(0.4), font_size=15, color=ACCENT, bold=True)

tbl_data = [
    ["기관", "주체", "핵심 임무"],
    ["시설", "요양(병)원", "화재 상황 전파·초기 대응, 입소 환자 대피 유도, 119 신고"],
    ["소방", "관할 소방서", "화재 진압, 현장 인명 구조·구급, 중증도 분류 및 병원 이송"],
    ["보건", "관할 보건소", "현장 응급의료 지원, 전원 환자 건강상태 파악"],
    ["지자체\n(기초)", "시·군·구", "현장 대응 인력 지원, 대규모 이송 차량 동원"],
    ["지자체\n(광역)", "시·도", "재난 상황 총괄 지휘·조정, 외부 임시 대피소 지정"],
]
add_table(slide, 6, 3, tbl_data,
          Inches(0.8), Inches(2.05), Inches(11.7), Inches(2.8),
          col_widths=[1.8, 2.2, 7.7])

# 3종 훈련 체계
add_title_text(slide, "현장 적용성을 높인 3종 훈련 체계",
               Inches(0.8), Inches(5.1), Inches(5), Inches(0.4), font_size=15, color=ACCENT, bold=True)

tbl_data2 = [
    ["훈련 유형", "방식", "활용"],
    ["도상훈련", "도면·인력배치표 기반 토론형", "종사자별 임무·수평 피난 동선 점검"],
    ["현장 모의훈련", "실전 기동형\n(침대 이동, 방화문 폐쇄 등)", "대본형 가이드 제공,\n시설 자체 분기별 실시"],
    ["전문가 참여 훈련", "외부 유관기관\n직접 참여 종합 검증", "자동속보 연동, 119 환자 인계,\n대피소 이동 등 총합 검증"],
]
add_table(slide, 4, 3, tbl_data2,
          Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3),
          col_widths=[2.5, 4.0, 5.2])


# ============================================================
# 슬라이드 18: 시계열 6단계 타임라인
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 18, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업4. 시계열 6단계 타임라인 진행표",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 6단계 시각화
stages = ["발화", "감지", "전파", "대피", "확인", "수습"]
stage_width = 1.7
start_x = 1.0
y_stage = 2.0

for i, stage in enumerate(stages):
    x = start_x + i * stage_width + i * 0.25
    # 단계 박스
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y_stage), Inches(stage_width), Inches(0.55))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY if i % 2 == 0 else DARK_BLUE
    box.line.fill.background()
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"STEP {i+1}. {stage}"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(4)
    
    # 화살표 (마지막 제외)
    if i < len(stages) - 1:
        arrow_x = x + stage_width
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            Inches(arrow_x), Inches(y_stage + 0.12), Inches(0.25), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GRAY
        arrow.line.fill.background()

# 주간/야간 비교
add_title_text(slide, "주간/야간 분리 시나리오",
               Inches(0.8), Inches(3.0), Inches(5), Inches(0.4), font_size=15, color=ACCENT, bold=True)

# 주간 박스
box_day = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(3.5), Inches(5.6), Inches(2.5))
box_day.fill.solid()
box_day.fill.fore_color.rgb = LIGHT_BG
box_day.line.color.rgb = DARK_BLUE
box_day.line.width = Pt(1)

add_title_text(slide, "☀  주간 시나리오",
               Inches(1.0), Inches(3.6), Inches(5.2), Inches(0.4),
               font_size=14, color=NAVY, bold=True)
add_bullet_list(slide, [
    "충분한 인력 확보 상태",
    "C유형 1차 수평 대피 실시",
    "상황 따라 수직 피난 전개 가능",
    "'적극적 대피 및 초기 진화' 중심",
], Inches(1.0), Inches(4.1), Inches(5.0), Inches(1.8), font_size=11)

# 야간 박스
box_night = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.9), Inches(3.5), Inches(5.6), Inches(2.5))
box_night.fill.solid()
box_night.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
box_night.line.color.rgb = NAVY
box_night.line.width = Pt(1)

add_title_text(slide, "🌙  야간·심야 시나리오",
               Inches(7.1), Inches(3.6), Inches(5.2), Inches(0.4),
               font_size=14, color=WHITE, bold=True)
add_bullet_list(slide, [
    "최소 인력 근무 상태",
    "외부 대피·수직 이동 전면 배제",
    "자동속보 119 전파 집중",
    "'수평 피난 + 생존 대기' 집중",
], Inches(7.1), Inches(4.1), Inches(5.0), Inches(1.8), font_size=11, color=WHITE, bullet_color=RGBColor(0xA0, 0xC4, 0xE8))


# ============================================================
# 슬라이드 19: 정책 과제 - 법령 개선
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 19, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업4. 정책 과제 ─ 법령 및 제도 개선",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 문제 제기
add_title_text(slide, "현행 법령의 한계",
               Inches(0.8), Inches(1.7), Inches(5), Inches(0.4), font_size=15, color=RED_ACCENT, bold=True)
add_bullet_list(slide, [
    "소관 부처별 용어 파편화: 요양병원(의료법), 노인요양시설(노인복지법), 노유자시설(소방시설법)",
    "재실자의 위험도(와상환자)는 동일하나 바닥 면적에 따라 설비 기준 제각각",
    "위락시설 등 화재 고위험 시설과의 용도 혼재 원천 금지 미비",
], Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.6), font_size=11)

# 개선 제안
add_title_text(slide, "개선 제안: 「피난약자시설」 단일 법적 카테고리 신설",
               Inches(0.8), Inches(3.7), Inches(10), Inches(0.4), font_size=15, color=DARK_BLUE, bold=True)

add_highlight_box(slide,
    "24시간 피난약자가 숙식하는 모든 시설(요양병원, 노인요양시설, 장애인 거주시설 등)을\n「피난약자시설」이라는 새로운 단일 법적 카테고리로 통합 신설",
    Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.9),
    bg_color=RGBColor(0xE8, 0xF4, 0xFD), text_color=NAVY, font_size=12, bold=True)

add_title_text(slide, "기대 효과",
               Inches(0.8), Inches(5.3), Inches(5), Inches(0.4), font_size=15, color=ACCENT, bold=True)
add_bullet_list(slide, [
    "시설 명칭·면적 관계없이 '스프링클러 전면 설치' 등 동일 안전 기준 예외 없이 적용",
    "직관적인 화재 안전망 구축으로 소규모 시설 사각지대 완전 해소",
], Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.9), font_size=11)


# ============================================================
# 슬라이드 20: A~D등급 평가 및 인센티브
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 20, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "과업4. 시설 자체 평가 도구 및 인센티브 제도",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# A~D 등급
add_title_text(slide, "A~D등급 체계 기반 시설 자체 평가 도구",
               Inches(0.8), Inches(1.65), Inches(8), Inches(0.4), font_size=15, color=ACCENT, bold=True)

add_body_text(slide, "4대 평가 계층: ①피난·방화계획(최우선) → ②화재안전관리·예방 → ③소방시설 점검 → ④소방대 활동 가능성",
              Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.4),
              font_size=11, color=GRAY)

tbl_data = [
    ["등급", "기준", "활용"],
    ["A", "법적 기준 \"초과\" + 추가 안전대책 보유", "우수시설 인증마크 대상"],
    ["B", "법적 기준 \"충족\"", "현행 유지"],
    ["C", "유지·관리 \"미흡\"", "방문 컨설팅 우선 대상"],
    ["D", "법적 기준 \"미달\"", "긴급 개선 대상"],
]
add_table(slide, 5, 3, tbl_data,
          Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0),
          col_widths=[1.5, 5.5, 4.7])

# 인센티브
add_title_text(slide, "민간 자발적 안전 투자 유도 인센티브 제도",
               Inches(0.8), Inches(4.8), Inches(8), Inches(0.4), font_size=15, color=ACCENT, bold=True)

add_bullet_list(slide, [
    "A등급 시설 대상 '안전관리 우수시설 인증마크' 부여 → 환자 유치·홍보 활용",
    "국민건강보험공단 요양기관 평가 및 지자체 재난관리평가 시 가점 부여",
    "민간 운영자의 자발적 소방시설 투자·훈련 참여를 유도하는 선순환 구조 구축",
], Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), font_size=12)


# ============================================================
# 슬라이드 21: 투-트랙 대피 전략
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 21, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "투-트랙(Two-Track) 맞춤형 대피 표준모델",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# Track 1 박스
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
box1.line.color.rgb = DARK_BLUE
box1.line.width = Pt(1.5)

add_title_text(slide, "Track 1", Inches(1.0), Inches(1.85), Inches(5.2), Inches(0.4),
               font_size=20, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "대형 요양병원", Inches(1.0), Inches(2.25), Inches(5.2), Inches(0.4),
               font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "수평 피난 및 제자리 방어 (Defend-in-Place)", Inches(1.0), Inches(2.65), Inches(5.2), Inches(0.3),
               font_size=11, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

add_divider_line(slide, Inches(1.2), Inches(3.05), Inches(5.0))

add_title_text(slide, "하드웨어적 대책", Inches(1.0), Inches(3.15), Inches(5.2), Inches(0.3),
               font_size=12, color=NAVY, bold=True)
add_bullet_list(slide, [
    "층별 최소 2개 이상 수평 방화구획 설치",
    "침대 진입 가능 피난용 승강기 확보",
    "열감지기 → 연기감지기 전면 교체",
    "배연창(감지기 연동) 설치",
], Inches(1.0), Inches(3.45), Inches(5.2), Inches(1.5), font_size=10)

add_title_text(slide, "소프트웨어적 대책 (행동 매뉴얼)", Inches(1.0), Inches(4.95), Inches(5.2), Inches(0.3),
               font_size=12, color=NAVY, bold=True)
add_bullet_list(slide, [
    "C→B→A유형 순 인접 방화구획 수평 이동",
    "야간: 수직 이동 배제, 수평피난+생존대기",
    "종사자 4단계 임무 규격화",
], Inches(1.0), Inches(5.25), Inches(5.2), Inches(1.2), font_size=10)

# Track 2 박스
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.9), Inches(1.7), Inches(5.6), Inches(5.0))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
box2.line.color.rgb = ORANGE
box2.line.width = Pt(1.5)

add_title_text(slide, "Track 2", Inches(7.1), Inches(1.85), Inches(5.2), Inches(0.4),
               font_size=20, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "소규모 요양원", Inches(7.1), Inches(2.25), Inches(5.2), Inches(0.4),
               font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "신속 옥외 피난", Inches(7.1), Inches(2.65), Inches(5.2), Inches(0.3),
               font_size=11, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

add_divider_line(slide, Inches(7.3), Inches(3.05), Inches(5.0))

add_title_text(slide, "하드웨어적 대책", Inches(7.1), Inches(3.15), Inches(5.2), Inches(0.3),
               font_size=12, color=NAVY, bold=True)
add_bullet_list(slide, [
    "단거리 옥외 대피용 경사로 확보",
    "넓은 출입구 확보",
    "자동화재속보설비 전면 완비",
], Inches(7.1), Inches(3.45), Inches(5.2), Inches(1.2), font_size=10)

add_title_text(slide, "소프트웨어적 대책 (행동 매뉴얼)", Inches(7.1), Inches(4.95), Inches(5.2), Inches(0.3),
               font_size=12, color=NAVY, bold=True)
add_bullet_list(slide, [
    "신속히 건물 밖 옥외 대피소로 이동",
    "야간: 초기 진화 미집착, 즉시 119 신고",
    "소수 인력 총력 환자 탈출 최우선",
], Inches(7.1), Inches(5.25), Inches(5.2), Inches(1.2), font_size=10)


# ============================================================
# 슬라이드 22: 중장기 확산 로드맵
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 22, TOTAL_SLIDES)
add_section_header(slide, "Ⅱ. 세부 연구 수행 계획")

add_title_text(slide, "중장기 확산 로드맵",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

# 3단계 로드맵
roadmap = [
    ("1단계", "도입·검증", "연구 완료 후\n6개월 이내", [
        "대표 시설 3~5개소 시범사업 운영",
        "유형별 대피계획 표준안 + 합동 훈련 시나리오 적용",
        "종사자 매뉴얼 이해도, 수평 피난 동선 현장 작동성 교차 검증",
    ]),
    ("2단계", "보급·교육", "모델 확정 후\n1~2년차", [
        "행안부·복지부·소방청·대한요양병원협회 간 공조 체계 구축",
        "고위험 시설 우선순위 선정 → 전문가 현장 방문 맞춤형 컨설팅",
        "전국 매뉴얼 배포",
    ]),
    ("3단계", "정착·제도화", "모델 확정 후\n2~3년차", [
        "법정 소방계획서 양식에 대피체계 표준안 통합",
        "합동 모의훈련 정례화, 지자체 정기 안전 점검 항목 반영",
        "A~D등급 자가진단 도구 전국 보급 → 상시 사후관리 체계 구축",
    ]),
]

x_positions = [0.8, 4.8, 8.8]
for i, (stage, name, period, items) in enumerate(roadmap):
    x = x_positions[i]
    # 배경 박스
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(1.7), Inches(3.7), Inches(5.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = RGBColor(0xD0, 0xDA, 0xE8)
    box.line.width = Pt(0.75)
    
    # 단계 배지
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.2), Inches(1.85), Inches(3.3), Inches(0.45))
    badge.fill.solid()
    colors = [DARK_BLUE, ACCENT, NAVY]
    badge.fill.fore_color.rgb = colors[i]
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = f"{stage}: {name}"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(3)
    
    # 기간
    add_body_text(slide, period, Inches(x + 0.2), Inches(2.45), Inches(3.3), Inches(0.5),
                  font_size=11, color=GRAY, bold=False)
    
    # 내용
    add_bullet_list(slide, items, Inches(x + 0.2), Inches(3.1), Inches(3.3), Inches(3.3),
                    font_size=10, bullet_color=colors[i])


# ============================================================
# 슬라이드 23: 사업관리 - 연구 추진 일정
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 23, TOTAL_SLIDES)
add_section_header(slide, "Ⅲ. 사업 관리 및 기관 역량")

add_title_text(slide, "연구 추진 일정 및 마일스톤",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

tbl_data = [
    ["단계", "수행 기간", "주요 과업", "핵심 산출물"],
    ["1단계", "1~2개월", "기초현황 조사,\n화재 대피 사례 분석,\n관계자 심층 인터뷰", "현황보고서\n유형분류표"],
    ["2단계", "2~4개월", "국내외 비교분석,\n유형별 대피체계\n표준모델 초안 개발", "비교분석보고서\n표준안 초안"],
    ["3단계", "4~5개월", "행동 매뉴얼,\n합동 훈련 시나리오,\n정책과제 발굴", "행동 매뉴얼\n훈련 시나리오\n정책과제 제안서"],
    ["4단계", "5~6개월", "전문가 자문·현장 검토,\n최종 모델 보완\n완료 보고회", "최종보고서\n요약보고서"],
]
add_table(slide, 5, 4, tbl_data,
          Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.5),
          col_widths=[1.5, 1.8, 4.5, 3.9])


# ============================================================
# 슬라이드 24: 보고 체계
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, WHITE)
add_top_accent_line(slide)
add_bottom_bar(slide)
add_page_number(slide, 24, TOTAL_SLIDES)
add_section_header(slide, "Ⅲ. 사업 관리 및 기관 역량")

add_title_text(slide, "보고 체계 및 품질관리",
               Inches(0.8), Inches(0.85), Inches(10), Inches(0.6), font_size=22, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.45), Inches(11.7))

tbl_data = [
    ["보고 유형", "시점", "주요 내용"],
    ["착수보고", "계약 후 10일 이내", "이력서, 보안대책, 용역비 산출내역, 추진일정·방향"],
    ["중간보고", "3~4개월차", "과업 진척상황, 중간결과 공유, 향후 계획 협의"],
    ["완료보고", "계약만료 10일 전", "최종 보고회 개최, 산출물 납품"],
    ["월간보고", "매월", "월간 공정 진행상황 서면 보고"],
    ["수시보고", "필요시", "행안부 요구 시 연구진행 상황 보고"],
]
add_table(slide, 6, 3, tbl_data,
          Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.5),
          col_widths=[2.0, 3.0, 6.7])

# FDS 특허 적용 방안
add_title_text(slide, "FDS 자동화 특허 과업 적용 방안",
               Inches(0.8), Inches(4.6), Inches(8), Inches(0.4), font_size=15, color=ACCENT, bold=True)

add_highlight_box(slide,
    "요양병원 유형별 3차원 객체·개구부·측정 지점 자동 생성 → FDS 표준 텍스트 파일 자동 변환\n→ ASET 산출의 객관성·정밀도 대폭 향상 + 시뮬레이션 소요시간 획기적 단축",
    Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.9),
    bg_color=RGBColor(0xE8, 0xF4, 0xFD), text_color=NAVY, font_size=12, bold=False)


# ============================================================
# 슬라이드 25: 마무리
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, NAVY)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = DARK_BLUE
line.line.fill.background()

add_title_text(slide, "감사합니다",
               Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
               font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.03))
divider.fill.solid()
divider.fill.fore_color.rgb = DARK_BLUE
divider.line.fill.background()

add_title_text(slide, "요양(병)원 시설 유형별 대피체계 표준모델 개발",
               Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6),
               font_size=18, color=RGBColor(0xA0, 0xC4, 0xE8), bold=False, alignment=PP_ALIGN.CENTER)

add_title_text(slide, "요약제안서",
               Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.5),
               font_size=14, color=RGBColor(0x88, 0xAA, 0xCC), bold=False, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# 파일 저장
# ──────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "요양병원_대피체계_요약제안서.pptx")
prs.save(output_path)
print(f"[OK] PPT file created: {output_path}")
print(f"   Total {TOTAL_SLIDES} slides")
