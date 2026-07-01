import sys
import os

pptx_path = "D:\\2026년\\요양병원시설 대피체계\\요양시설_대피_표준모델.pptx"
txt_path = "D:\\2026년\\요양병원시설 대피체계\\요양시설_대피_표준모델.txt"

def extract_text():
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            text += f"--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(text)
        print("Successfully extracted using python-pptx")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_text()
